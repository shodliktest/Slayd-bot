"""
PowerPoint Presentation Generator
Creates professional presentations with images and styling
"""
import logging
import os
from typing import Dict, List
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from config import SLIDE_THEMES, DEFAULT_THEME, TEMP_FILES_DIR
from .web_search import find_images, download_image

logger = logging.getLogger(__name__)


async def generate_pptx(slides_data: Dict, output_path: str) -> str:
    """Generate PowerPoint presentation"""
    
    try:
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Get theme
        theme_name = slides_data.get('theme', DEFAULT_THEME)
        theme = SLIDE_THEMES.get(theme_name, SLIDE_THEMES[DEFAULT_THEME])
        
        # Title slide
        await _add_title_slide(prs, slides_data.get('title', 'Presentation'), theme)
        
        # Content slides
        slides = slides_data.get('slides', [])
        
        for slide_info in slides:
            await _add_content_slide(prs, slide_info, theme)
        
        # Save presentation
        prs.save(output_path)
        
        logger.info(f"Generated PowerPoint: {output_path} ({len(slides)} slides)")
        return output_path
        
    except Exception as e:
        logger.error(f"Error generating PowerPoint: {e}")
        raise


async def _add_title_slide(prs, title: str, theme: Dict):
    """Add title slide"""
    
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*theme['bg_color'])
    
    # Title
    left = Inches(1)
    top = Inches(3)
    width = Inches(8)
    height = Inches(1.5)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = title
    
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*theme['title_color'])
    p.alignment = PP_ALIGN.CENTER


async def _add_content_slide(prs, slide_info: Dict, theme: Dict):
    """Add content slide with text and optionally image"""
    
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*theme['bg_color'])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    tf.text = slide_info.get('title', '')
    
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*theme['title_color'])
    
    # Content
    content_items = slide_info.get('content', [])
    
    # Check if we need to add image
    image_keywords = slide_info.get('image_keywords', '')
    has_image = False
    image_path = None
    
    if image_keywords:
        try:
            images = await find_images(image_keywords, count=1)
            if images:
                image_path = os.path.join(TEMP_FILES_DIR, f"slide_{slide_info['slide_number']}.jpg")
                if await download_image(images[0], image_path):
                    has_image = True
        except Exception as e:
            logger.warning(f"Could not add image: {e}")
    
    # Layout: text on left, image on right (if available)
    if has_image and image_path:
        # Text box (left side)
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(5), Inches(5))
        
        # Image (right side)
        try:
            slide.shapes.add_picture(image_path, Inches(6), Inches(2), width=Inches(3.5))
        except Exception as e:
            logger.warning(f"Could not add image to slide: {e}")
    else:
        # Full width text box
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(5))
    
    # Add content bullets
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(*theme['text_color'])
        p.space_before = Pt(12)
        p.level = 0


async def create_presentation_from_data(
    content_id: str,
    slides_data: Dict
) -> str:
    """Create PowerPoint file from slides data"""
    
    output_path = os.path.join(TEMP_FILES_DIR, f"{content_id}.pptx")
    
    await generate_pptx(slides_data, output_path)
    
    return output_path
